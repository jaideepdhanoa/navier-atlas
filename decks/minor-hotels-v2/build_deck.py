#!/usr/bin/env python3
"""
Hospitality Deck v2 builder (reference instance: Minor Hotels x Navier).
Spec-driven so Grok can codify it into a generator. Reads a DECK spec (copy +
image bindings), bakes legibility scrims, and emits a 16:9 PPTX which is then
upload-converted to Google Slides. ALL copy is plain partner-facing English.
"""
import os, json, numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = "/tasklet/agent/home/hospitality-template-2026-06-25"
A = f"{ROOT}/assets"
S = f"{ROOT}/assets/scrim"
os.makedirs(S, exist_ok=True)

# ---------- palette ----------
INK   = RGBColor(0x0A,0x12,0x1C)
PAPER = RGBColor(0xF6,0xF3,0xEC)
GOLD  = RGBColor(0xC9,0xA3,0x5B)
MUTE  = RGBColor(0xC2,0xCD,0xD8)
SERIF = "Playfair Display"
SANS  = "Arial"
EW, EH = Inches(13.333), Inches(7.5)

# ---------- scrim baking ----------
def _bake(src, dst, amap):
    img = Image.open(src).convert("RGB")
    if img.size != (1536,864): img = img.resize((1536,864))
    arr = np.asarray(img).astype(float)
    dark = np.array([6,12,20], float)
    a = np.clip(amap,0,0.92)[...,None]
    out = arr*(1-a)+dark*a
    Image.fromarray(out.astype("uint8")).save(dst, quality=95)
    return dst

def left_map(frac=0.50, amax=0.90):
    w,h=1536,864
    x=np.linspace(0,1,w)
    a=np.clip((frac-x)/(frac*0.5),0,1)*amax
    return np.tile(a,(h,1))

def bottom_map(frac=0.55, amax=0.86):
    w,h=1536,864
    y=np.linspace(0,1,h)
    a=np.clip((y-(1-frac))/frac,0,1)*amax
    return np.tile(a[:,None],(1,w))

def overall_map(base=0.40, bot=0.40):
    return np.clip(np.full((864,1536),base)+bottom_map(0.6,bot),0,0.9)

def scrim(name, kind):
    src=f"{A}/{name}"; dst=f"{S}/{kind}-{name}"
    if kind=="left":    return _bake(src,dst,left_map())
    if kind=="bottom":  return _bake(src,dst,bottom_map())
    if kind=="overall": return _bake(src,dst,overall_map())
    return src

# ---------- text helpers ----------
def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide,color=INK):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=color

def full_image(slide,path):
    slide.shapes.add_picture(path,0,0,EW,EH)

def tb(slide,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    box=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=box.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    return tf

def para(tf,text,size,color,bold=False,font=SANS,first=False,align=PP_ALIGN.LEFT,
         before=0,after=6,spacing=1.0,tracking=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_before=Pt(before); p.space_after=Pt(after); p.line_spacing=spacing
    r=p.add_run(); r.text=text
    f=r.font; f.size=Pt(size); f.bold=bold; f.name=font; f.color.rgb=color
    return p

def rule(slide,l,t,w,color=GOLD,h=0.028):
    from pptx.enum.shapes import MSO_SHAPE
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=color; s.line.fill.background()
    return s

def eyebrow(tf,text,first=True):
    para(tf,text.upper(),12.5,GOLD,bold=True,font=SANS,first=first,after=10)

# ---------- slide renderers ----------
def r_cover(prs,d):
    s=slide_blank(prs); full_image(s,scrim(d["image"],"bottom"))
    tf=tb(s,0.85,3.7,11.6,3.4)
    eyebrow(tf,d["eyebrow"])
    for i,ln in enumerate(d["title"].split("\n")):
        para(tf,ln,52,PAPER,bold=True,font=SERIF,after=0,spacing=1.0)
    para(tf,d["subtitle"],17,MUTE,font=SANS,before=14,after=0,spacing=1.1)
    tf2=tb(s,0.85,6.85,11.6,0.5)
    para(tf2,d["footer"],10.5,MUTE,font=SANS,first=True)

def r_exec(prs,d):
    s=slide_blank(prs); full_image(s,scrim(d["image"],"overall"))
    tf=tb(s,0.85,0.6,11.6,1.1); eyebrow(tf,d["eyebrow"])
    para(tf,d["title"],34,PAPER,bold=True,font=SERIF,after=0)
    quads=d["quads"]; xs=[0.85,6.95]; ys=[2.15,4.75]; i=0
    for ry in ys:
        for rx in xs:
            q=quads[i]; i+=1
            qt=tb(s,rx,ry,5.55,2.4)
            para(qt,q["label"].upper(),12,GOLD,bold=True,font=SANS,first=True,after=5)
            para(qt,q["text"],13.5,PAPER,font=SANS,spacing=1.12)

def r_concept_center(prs,d):
    s=slide_blank(prs); full_image(s,scrim(d["image"],"overall"))
    tf=tb(s,0.85,2.5,11.6,2.5,anchor=MSO_ANCHOR.MIDDLE)
    eyebrow(tf,d["eyebrow"])
    para(tf,d["words"],46,PAPER,bold=True,font=SERIF,after=12)
    para(tf,d["body"],16,MUTE,font=SANS,spacing=1.15)

def r_bullets(prs,d):
    s=slide_blank(prs); fill_bg(s)
    if d.get("image"): 
        s.shapes.add_picture(scrim(d["image"],"left"),Inches(7.0),0,Inches(6.333),EH)
    tfe=tb(s,0.85,0.85,6.0,0.55)
    eyebrow(tfe,d["eyebrow"])
    rule(s,0.9,1.42,0.9)
    tf=tb(s,0.85,1.7,6.0,5.1)
    para(tf,d["title"],32,PAPER,bold=True,font=SERIF,first=True,after=8)
    for b in d["bullets"]:
        para(tf,b,14.5,MUTE,font=SANS,before=10,after=0,spacing=1.15)

def r_three(prs,d):
    s=slide_blank(prs); fill_bg(s)
    tf=tb(s,0.85,0.8,11.6,1.4); eyebrow(tf,d["eyebrow"])
    para(tf,d["title"],32,PAPER,bold=True,font=SERIF,after=0)
    cols=d["cols"]; w=3.7; gap=0.35; x0=0.85
    for i,c in enumerate(cols):
        x=x0+i*(w+gap)
        rule(s,x,2.55,0.65)
        ct=tb(s,x,2.75,w,3.8)
        para(ct,c["head"],19,GOLD,bold=True,font=SERIF,first=True,after=8)
        para(ct,c["text"],13.5,PAPER if False else MUTE,font=SANS,spacing=1.18)

def r_specs(prs,d):
    s=slide_blank(prs); fill_bg(s)
    tf=tb(s,0.85,0.8,11.6,1.4); eyebrow(tf,d["eyebrow"])
    para(tf,d["title"],32,PAPER,bold=True,font=SERIF,after=0)
    items=d["items"]; w=2.75; gap=0.32; x0=0.85
    for i,it in enumerate(items):
        x=x0+i*(w+gap)
        ct=tb(s,x,2.7,w,3.6)
        para(ct,it["stat"],38,GOLD,bold=True,font=SERIF,first=True,after=2)
        para(ct,it["unit"],12.5,PAPER,bold=True,font=SANS,after=8)
        para(ct,it["note"],12.5,MUTE,font=SANS,spacing=1.15)

def r_footprint(prs,d):
    s=slide_blank(prs); fill_bg(s)
    tf=tb(s,0.85,0.8,11.6,1.7); eyebrow(tf,d["eyebrow"])
    para(tf,d["title"],32,PAPER,bold=True,font=SERIF,after=8)
    para(tf,d["intro"],14.5,MUTE,font=SANS,spacing=1.15)
    rows=d["rows"]; y=3.15
    for rrow in rows:
        rt=tb(s,0.85,y,11.6,0.9)
        para(rt,rrow["name"],17,GOLD,bold=True,font=SERIF,first=True,after=1)
        para(rt,rrow["detail"],12.5,MUTE,font=SANS,spacing=1.05)
        rule(s,0.9,y+0.78,11.55,color=RGBColor(0x24,0x32,0x42),h=0.012)
        y+=0.82

def r_cluster(prs,d):
    s=slide_blank(prs); full_image(s,scrim(d["image"],"left"))
    tf=tb(s,0.8,0.75,5.9,6.3)
    eyebrow(tf,d["eyebrow"])
    para(tf,d["name"],33,PAPER,bold=True,font=SERIF,after=2)
    para(tf,d["subtitle"],12.5,GOLD,bold=True,font=SANS,after=12)
    for blk in d["blocks"]:
        para(tf,blk["label"].upper(),11,GOLD,bold=True,font=SANS,before=8,after=3)
        para(tf,blk["text"],12.5,PAPER,font=SANS,spacing=1.12)

def r_close(prs,d):
    s=slide_blank(prs); full_image(s,scrim(d["image"],"bottom"))
    tf=tb(s,0.85,3.9,11.6,3.0)
    for ln in d["title"].split("\n"):
        para(tf,ln,40,PAPER,bold=True,font=SERIF,after=0)
    para(tf,d["body"],15.5,MUTE,font=SANS,before=14,spacing=1.18)
    tf2=tb(s,0.85,6.9,11.6,0.45)
    para(tf2,d["footer"],10.5,MUTE,font=SANS,first=True)

def r_appendix(prs,d):
    s=slide_blank(prs); fill_bg(s)
    tf=tb(s,0.85,2.7,11.6,2.2,anchor=MSO_ANCHOR.MIDDLE)
    eyebrow(tf,d["eyebrow"])
    para(tf,d["title"],36,PAPER,bold=True,font=SERIF,after=12)
    para(tf,d["note"],14,MUTE,font=SANS,spacing=1.2)

def r_econ(prs,d):
    s=slide_blank(prs); fill_bg(s)
    tf=tb(s,0.85,0.8,11.6,1.7)
    eyebrow(tf,d["eyebrow"])
    para(tf,d["title"],27,PAPER,bold=True,font=SERIF,after=4)
    para(tf,d["corridor"],13,GOLD,bold=True,font=SANS,after=0)
    # left: assumptions
    at=tb(s,0.85,2.85,5.4,4.0)
    para(at,"WHAT WE ASSUME",11,GOLD,bold=True,font=SANS,first=True,after=8)
    for k,v in d["assumptions"]:
        p=at.add_paragraph(); p.space_after=Pt(7)
        r1=p.add_run(); r1.text=k+"   "; r1.font.size=Pt(13); r1.font.name=SANS; r1.font.color.rgb=MUTE
        r2=p.add_run(); r2.text=v; r2.font.size=Pt(13); r2.font.bold=True; r2.font.name=SANS; r2.font.color.rgb=PAPER
    # right: result card
    from pptx.enum.shapes import MSO_SHAPE
    card=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(6.7),Inches(2.7),Inches(5.75),Inches(3.7))
    card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0x10,0x1B,0x28); card.line.color.rgb=GOLD; card.line.width=Pt(1)
    rt=tb(s,7.05,3.0,5.1,3.2)
    para(rt,"THE RESULT, PER VESSEL / YEAR",11,GOLD,bold=True,font=SANS,first=True,after=10)
    for lbl,val,big in d["results"]:
        p=rt.add_paragraph(); p.space_after=Pt(6)
        r1=p.add_run(); r1.text=lbl+"  "; r1.font.size=Pt(13); r1.font.name=SANS; r1.font.color.rgb=MUTE
        r2=p.add_run(); r2.text=val; r2.font.size=Pt(22 if big else 14); r2.font.bold=True
        r2.font.name=SERIF if big else SANS; r2.font.color.rgb=(GOLD if big else PAPER)
    ft=tb(s,0.85,6.75,11.6,0.5)
    para(ft,d["foot"],9.5,MUTE,font=SANS,first=True)

RENDER={"cover":r_cover,"exec":r_exec,"concept":r_concept_center,"bullets":r_bullets,
        "three":r_three,"specs":r_specs,"footprint":r_footprint,"cluster":r_cluster,
        "close":r_close,"appendix":r_appendix,"econ":r_econ}

def build(spec, out):
    prs=Presentation(); prs.slide_width=EW; prs.slide_height=EH
    for sl in spec["slides"]:
        RENDER[sl["type"]](prs,sl)
    prs.save(out)
    print("saved",out,"slides:",len(spec["slides"]))

if __name__=="__main__":
    spec=json.load(open(f"{ROOT}/deck_spec_minor.json"))
    build(spec,f"{ROOT}/Minor-Hotels-Navier-v2.pptx")
